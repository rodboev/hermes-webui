from __future__ import annotations

import io
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from openpyxl import load_workbook
from pptx import Presentation

CLAIMED_OFFICE_EXTENSIONS = frozenset({".docx", ".xlsx", ".pptx"})
CLAIMED_OFFICE_FORMATS = frozenset({"docx", "xlsx", "pptx"})
OFFICE_PREVIEW_KIND = "office"
OFFICE_RENDER_MODE = "code"

_DOCX_BODY_CHILDREN = {qn("w:p"), qn("w:sectPr")}
_DOCX_PARAGRAPH_CHILDREN = {qn("w:r")}
_DOCX_RUN_CHILDREN = {qn("w:rPr"), qn("w:t")}


def is_claimed_office_path(path: str | Path) -> bool:
    return Path(str(path)).suffix.lower() in CLAIMED_OFFICE_EXTENSIONS


def _office_format_for_path(path: str | Path) -> str:
    return Path(str(path)).suffix.lower().lstrip(".")


def _normalise_preview_text(value) -> str:
    if value is None:
        return ""
    return str(value).replace("\r", "\n").replace("\n", " ").strip()


def _preview_line_count(content: str) -> int:
    if not content:
        return 1
    return content.count("\n") + 1


def _docx_preview_text(document: Document) -> str:
    chunks: list[str] = []
    for paragraph in document.paragraphs:
        chunks.append(paragraph.text or "")
    for table_index, table in enumerate(document.tables, start=1):
        table_lines = [f"Table {table_index}"]
        for row in table.rows:
            cells = [_normalise_preview_text(cell.text) for cell in row.cells]
            table_lines.append("\t".join(cells))
        chunks.append("\n".join(table_lines))
    return "\n\n".join(chunks).strip()


def _docx_editability(document: Document) -> tuple[bool, str | None]:
    body = document._element.body
    for child in body:
        if child.tag not in _DOCX_BODY_CHILDREN:
            return False, "docx contains unsupported structures"
    for paragraph in document.paragraphs:
        for child in paragraph._p:
            if child.tag not in _DOCX_PARAGRAPH_CHILDREN:
                return False, "docx contains unsupported paragraph structures"
        for run in paragraph.runs:
            for child in run._r:
                if child.tag not in _DOCX_RUN_CHILDREN:
                    return False, "docx contains unsupported inline content"
    return True, None


def _preview_docx(raw: bytes) -> tuple[str, bool, str | None]:
    try:
        document = Document(io.BytesIO(raw))
    except Exception as exc:  # pragma: no cover - library-specific failure mode
        raise ValueError("Unable to read DOCX preview") from exc
    content = _docx_preview_text(document)
    editable, reason = _docx_editability(document)
    return content, editable, reason


def _preview_xlsx(raw: bytes) -> str:
    try:
        workbook = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    except Exception as exc:  # pragma: no cover - library-specific failure mode
        raise ValueError("Unable to read XLSX preview") from exc
    chunks: list[str] = []
    for sheet in workbook.worksheets:
        sheet_lines = [f"Sheet: {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            values = [_normalise_preview_text(value) for value in row]
            if any(values):
                sheet_lines.append("\t".join(values))
        chunks.append("\n".join(sheet_lines).strip())
    if not chunks:
        return "Empty workbook"
    return "\n\n".join(chunk for chunk in chunks if chunk).strip() or "Empty workbook"


def _preview_pptx(raw: bytes) -> str:
    try:
        presentation = Presentation(io.BytesIO(raw))
    except Exception as exc:  # pragma: no cover - library-specific failure mode
        raise ValueError("Unable to read PPTX preview") from exc
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_lines = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            text = _normalise_preview_text(getattr(shape, "text", ""))
            if text:
                slide_lines.append(text)
        if len(slide_lines) == 1:
            slide_lines.append("(empty slide)")
        chunks.append("\n".join(slide_lines).strip())
    if not chunks:
        return "Empty presentation"
    return "\n\n".join(chunk for chunk in chunks if chunk).strip() or "Empty presentation"


def preview_office_document(path: str | Path, raw: bytes) -> dict:
    office_format = _office_format_for_path(path)
    if office_format not in CLAIMED_OFFICE_FORMATS:
        raise ValueError(f"Unsupported Office format: {path}")

    if office_format == "docx":
        content, editable, reason = _preview_docx(raw)
    elif office_format == "xlsx":
        content = _preview_xlsx(raw)
        editable, reason = False, "xlsx preview is read-only in this slice"
    elif office_format == "pptx":
        content = _preview_pptx(raw)
        editable, reason = False, "pptx preview is read-only in this slice"
    else:  # pragma: no cover - exhaustive guard
        raise ValueError(f"Unsupported Office format: {path}")

    payload = {
        "path": str(path),
        "content": content,
        "size": len(raw),
        "lines": _preview_line_count(content),
        "preview_kind": OFFICE_PREVIEW_KIND,
        "office_format": office_format,
        "render_mode": OFFICE_RENDER_MODE,
        "editable": editable,
    }
    if reason:
        payload["edit_blocked_reason"] = reason
    return payload


def _docx_bytes_from_text(content: str) -> bytes:
    document = Document()
    body = document._element.body
    for child in list(body):
        if child.tag != qn("w:sectPr"):
            body.remove(child)
    text = str(content or "").replace("\r\n", "\n").replace("\r", "\n")
    for line in text.split("\n"):
        document.add_paragraph(line)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def save_office_document(path: str | Path, current_bytes: bytes, content: str) -> tuple[dict, bytes]:
    office_format = _office_format_for_path(path)
    if office_format != "docx":
        raise ValueError(f"{office_format or 'office file'} is preview-only in this slice")

    current_preview = preview_office_document(path, current_bytes)
    if not current_preview.get("editable"):
        raise ValueError(current_preview.get("edit_blocked_reason") or "DOCX document is not editable")

    saved_bytes = _docx_bytes_from_text(content)
    saved_preview = preview_office_document(path, saved_bytes)
    if not saved_preview.get("editable"):
        raise ValueError(saved_preview.get("edit_blocked_reason") or "Saved DOCX is not editable")
    return saved_preview, saved_bytes
